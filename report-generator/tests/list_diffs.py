"""
Helper functions for testing differences between two generated reports for comparing versions
For now only compares text
"""

#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import os
from difflib import Differ

from pptx import Presentation


def compare_pptx(file1, file2):
    prs1 = Presentation(file1)
    prs2 = Presentation(file2)

    differences = []

    # Compare slide count
    if len(prs1.slides) != len(prs2.slides):
        differences.append(
            f"Slide count mismatch: {file1} has {len(prs1.slides)} slides, {file2} has {len(prs2.slides)} slides")

    # Compare slides
    for i, (slide1, slide2) in enumerate(zip(prs1.slides, prs2.slides)):
        # Compare shapes count
        if len(slide1.shapes) != len(slide2.shapes):
            differences.append(
                f"Slide {i + 1}: Shape count mismatch: {file1} has {len(slide1.shapes)} shapes, {file2} has {len(slide2.shapes)} shapes")

        # Compare text in shapes
        for j, (shape1, shape2) in enumerate(zip(slide1.shapes, slide2.shapes)):
            if hasattr(shape1, 'text') and hasattr(shape2, 'text'):
                if shape1.text != shape2.text:
                    differ = Differ()
                    diff = list(differ.compare(shape1.text.splitlines(), shape2.text.splitlines()))
                    differences.append(f"Slide {i + 1}, Shape {j + 1}: Text difference:")
                    differences.extend(
                        [f"    {line}" for line in diff if line.startswith('+ ') or line.startswith('- ')])

    return differences


def main():
    file1 = "../new.pptx"
    file2 = "../old.pptx"

    if not os.path.exists(file1) or not os.path.exists(file2):
        print("One or both files do not exist.")
        return

    differences = compare_pptx(file1, file2)

    if differences:
        print("Differences found:")
        for diff in differences:
            print(diff)
    else:
        print("No differences found.")


if __name__ == "__main__":
    main()
