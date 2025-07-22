#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import os
from difflib import Differ

from pptx import Presentation


def compare_pptx(file1, file2):
    prs1 = Presentation(file1)
    prs2 = Presentation(file2)

    differences = []

    differences += compare_slide_count(prs1, prs2, file1, file2)
    differences += compare_slides(prs1, prs2, file1, file2)

    return not differences, differences


def compare_slide_count(prs1, prs2, file1, file2):
    differences = []
    if len(prs1.slides) != len(prs2.slides):
        differences.append(
            f"Slide count mismatch: {file1} has {len(prs1.slides)} slides, {file2} has {len(prs2.slides)} slides"
        )
    return differences


def compare_slides(prs1, prs2, file1, file2):
    differences = []
    for i, (slide1, slide2) in enumerate(zip(prs1.slides, prs2.slides)):
        differences += compare_shape_count(slide1, slide2, file1, file2, i)
        differences += compare_shapes_text(slide1, slide2, i)
        differences += compare_tables(slide1, slide2, file1, file2, i)
    return differences


def compare_shape_count(slide1, slide2, file1, file2, slide_index):
    differences = []
    if len(slide1.shapes) != len(slide2.shapes):
        differences.append(
            f"Slide {slide_index + 1}: Shape count mismatch: "
            f"{file1} has {len(slide1.shapes)} shapes, {file2} has {len(slide2.shapes)} shapes"
        )
    return differences


def compare_shapes_text(slide1, slide2, slide_index):
    differences = []
    for j, (shape1, shape2) in enumerate(zip(slide1.shapes, slide2.shapes)):
        if shape1.has_text_frame and shape2.has_text_frame:
            if shape1.text != shape2.text:
                differ = Differ()
                diff = list(differ.compare(shape1.text.splitlines(), shape2.text.splitlines()))
                differences.append(f"Slide {slide_index + 1}, Shape {j + 1}: Text difference:")
                differences.extend(
                    [f"    {line}" for line in diff if line.startswith('+ ') or line.startswith('- ')]
                )
    return differences


def compare_tables(slide1, slide2, file1, file2, slide_index):
    differences = []
    for j, (shape1, shape2) in enumerate(zip(slide1.shapes, slide2.shapes)):
        if shape1.has_table and shape2.has_table:
            table1, table2 = shape1.table, shape2.table
            if len(table1.rows) != len(table2.rows) or len(table1.columns) != len(table2.columns):
                differences.append(
                    f"Slide {slide_index + 1}, Table {j + 1}: Table size mismatch: "
                    f"{file1} has {len(table1.rows)} rows and {len(table1.columns)} columns, "
                    f"{file2} has {len(table2.rows)} rows and {len(table2.columns)} columns"
                )
            else:
                differences += compare_table_cells(table1, table2, slide_index, j)
    return differences


def compare_table_cells(table1, table2, slide_index, table_index):
    differences = []
    for row_idx, (row1, row2) in enumerate(zip(table1.rows, table2.rows)):
        for col_idx, (cell1, cell2) in enumerate(zip(row1.cells, row2.cells)):
            if cell1.text != cell2.text:
                differ = Differ()
                diff = list(differ.compare(cell1.text.splitlines(), cell2.text.splitlines()))
                differences.append(
                    f"Slide {slide_index + 1}, Table {table_index + 1}, Cell ({row_idx + 1}, {col_idx + 1}): Text difference:"
                )
                differences.extend(
                    [f"    {line}" for line in diff if line.startswith('+ ') or line.startswith('- ')]
                )
    return differences


def main():
    file1 = "reference_output.pptx"
    file2 = "test_output.pptx"

    if not os.path.exists(file1) or not os.path.exists(file2):
        print("One or both files do not exist.")
        return

    are_equal, differences = compare_pptx(file1, file2)

    if not are_equal:
        print("Differences found:")
        for diff in differences:
            print(diff)
    else:
        print("No differences found.")


if __name__ == "__main__":
    main()
