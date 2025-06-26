#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from abc import ABC, abstractmethod
from typing import Callable

from pptx import Presentation
from pptx.chart.data import CategoryChartData

from report_generator.generator import report_utils
from report_generator.generator.data_models import maintainability_data, modernization_data
from report_generator.generator.placeholders import Placeholder
from report_generator.generator.placeholders.base import PlaceholderDocType


class _AbstractCategoryChartPlaceholder(Placeholder, ABC):
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    @abstractmethod
    def labels(cls):
        pass

    @classmethod
    @abstractmethod
    def series(cls):
        pass

    @classmethod
    def colors(cls):
        return []

    @classmethod
    @abstractmethod
    def axis_label(cls):
        pass

    @classmethod
    def value(cls, placeholder=None):
        return {
            "labels"   : cls.labels(),
            "series"   : cls.series(),
            "colors"   : cls.colors(),
            "axisLabel": cls.axis_label(),
        }

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable):
        charts = []
        for slide in report_utils.pptx.identify_specific_slide(presentation, key):
            for shape in slide.shapes:
                if shape.has_chart:
                    charts.append(shape.chart)

        if len(charts) == 0:
            return

        values = value_cb()
        chart_data = CategoryChartData()
        chart_data.categories = values["labels"]
        for y in values["series"]:
            chart_data.add_series(values["axisLabel"], y)

        colors = values["colors"]
        for chart in charts:
            chart.replace_data(chart_data)

            if not colors:
                continue

            for serie in chart.series:
                for idx, point in enumerate(serie.points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = colors[idx]


class TechnologyCategoryChartPlaceholder(_AbstractCategoryChartPlaceholder):
    """Chart with volume (in % of person months of code) per technology."""
    key = "TECHNOLOGY_CHART"

    @classmethod
    def labels(cls):
        return [data["displayName"] for data in maintainability_data.sorted_tech]

    @classmethod
    def series(cls):
        return [[data["volumeInPersonMonths"] / maintainability_data.tech_total_volume_pm for data in
                 maintainability_data.sorted_tech]]

    @classmethod
    def axis_label(cls):
        return "Volume in Person Months"


class TestCodeRatioCategoryChartPlaceholder(_AbstractCategoryChartPlaceholder):
    """Pie chart with volume and % of test code per technology, colored in line with the SIG test code benchmark."""
    key = "TEST_CODE_RATIO_CHART"

    @classmethod
    def labels(cls):
        return [data["displayName"] for data in maintainability_data.sorted_tech]

    @classmethod
    def series(cls):
        return [[data["volumeInPersonMonths"] / maintainability_data.tech_total_volume_pm for data in
                 maintainability_data.sorted_tech]]

    @classmethod
    def colors(cls):
        return [report_utils.pptx.test_code_ratio_color(data["testCodeRatio"]) for data in
                maintainability_data.sorted_tech]

    @classmethod
    def axis_label(cls):
        return "Volume in Person Months"


class TechnicalDebtSystemsChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "TECHNICAL_DEBT_SYSTEMS_CHART"

    @classmethod
    def labels(cls):
        candidates = modernization_data.modernization_candidates_by_technical_debt
        return [candidate.display_name for candidate in candidates]

    @classmethod
    def series(cls):
        candidates = modernization_data.modernization_candidates_by_technical_debt
        technical_debt = [candidate.technical_debt_in_py for candidate in candidates]
        remaining = [candidate.volume_in_py - candidate.technical_debt_in_py for candidate in candidates]
        return [technical_debt, remaining]

    @classmethod
    def axis_label(cls):
        return "Code volume in person years"
