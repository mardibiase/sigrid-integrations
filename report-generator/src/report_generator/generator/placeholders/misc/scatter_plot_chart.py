#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor

from report_generator.generator import report_utils
from report_generator.generator.data_models import modernization_data
from report_generator.generator.placeholders import Placeholder
from report_generator.generator.placeholders.base import PlaceholderDocType, PlaceholderDocsMetadata

BLUE_GRADIENT = ["003DAB", "2E6BFF", "8DA8FF", "DBE1FF", "8A98A8"]


class ModernizationScatterPlotChartPlaceholder(Placeholder):
    key = "MODERNIZATION_SCATTER_PLOT_CHART"
    __doc_metadata__ = PlaceholderDocsMetadata(type=PlaceholderDocType.CHART)

    @classmethod
    def value(cls, parameter=None):
        pass

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb) -> None:
        charts = report_utils.pptx.gather_charts(presentation, key)

        if len(charts) == 0:
            return

        grouped_candidates = {criticality: [] for criticality in ["CRITICAL", "HIGH", "MEDIUM", "LOW", "UNKNOWN"]}
        for candidate in modernization_data.modernization_candidates:
            grouped_candidates[candidate.business_criticality.upper()].append(candidate)

        chart_data = BubbleChartData()
        for group, candidates in grouped_candidates.items():
            series = chart_data.add_series(group.title())
            for c in candidates:
                series.add_data_point(c.estimated_effort_py, c.estimated_change_speed, c.technical_debt_in_py)

        for chart in charts:
            chart.replace_data(chart_data)
            for i, group in enumerate(grouped_candidates):
                for j, candidate in enumerate(grouped_candidates[group]):
                    chart.series[i].points[j].data_label.text_frame.text = candidate.display_name
                    chart.series[i].points[j].format.line.color.rgb = RGBColor(255, 255, 255)
                    chart.series[i].points[j].format.fill.solid()
                    chart.series[i].points[j].format.fill.fore_color.rgb = RGBColor.from_string(BLUE_GRADIENT[i])
            chart.value_axis.minimum_scale = 0
            chart.category_axis.minimum_scale = 0
