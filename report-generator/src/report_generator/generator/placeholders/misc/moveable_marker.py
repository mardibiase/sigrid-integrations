#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from abc import ABC
from typing import Callable

from pptx import Presentation

from report_generator.generator import report_utils
from report_generator.generator.data_models import architecture_data, maintainability_data, modernization_data, osh_data
from report_generator.generator.formatters.formatters import maintainability_round
from report_generator.generator.placeholders import Placeholder

_RATING_MARKER_MOVE_SIZE = 2200000
_MANAGEMENT_SUMMARY_MARKER_RANGE = 4000000


def _distance_to_average(rating):
    return -1 * (3.0 - rating)


class _AbstractMoveableMarkerPlaceholder(Placeholder, ABC):
    """Used for the rating indicator on SIG metric rating slides. It adds the rating, but also moves the marker to the correct position on the slide."""

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable[[], str]) -> None:
        paragraphs = []
        for slide in presentation.slides:
            paragraphs.extend(report_utils.pptx.find_text_in_slide(slide, key))

        if len(paragraphs) == 0:
            return

        value = value_cb()
        value_float = float(value)

        for paragraph in paragraphs:
            report_utils.pptx.update_paragraph(paragraph, key, value)

            # noinspection PyProtectedMember
            marker = paragraph._parent._parent._parent._parent
            marker.left = int(marker.left + _RATING_MARKER_MOVE_SIZE * _distance_to_average(value_float))


class MaintainabilityMovableMarkerPlaceholder(_AbstractMoveableMarkerPlaceholder):
    key = "MARKER_MAINT_RATING"

    @classmethod
    def value(cls, parameter=None):
        return maintainability_round(maintainability_data.maintainability_rating)


class ArchitectureMovableMarkerPlaceholder(_AbstractMoveableMarkerPlaceholder):
    key = "MARKER_ARCH_RATING"

    @classmethod
    def value(cls, parameter=None):
        return maintainability_round(architecture_data.ratings["architecture"])


class OSHMovableMarkerPlaceholder(_AbstractMoveableMarkerPlaceholder):
    key = "MARKER_OSH_RATING"

    @classmethod
    def value(cls, parameter=None) -> tuple[float, str]:
        return maintainability_round(osh_data.data.ratings["system"])


class _ManagementSummaryMarkerPlaceholder(Placeholder, ABC):
    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable[[], str]) -> None:
        for slide in presentation.slides:
            for marker in report_utils.pptx.find_text_in_slide(slide, key):
                value, label = value_cb()
                report_utils.pptx.update_paragraph(marker, key, f"{label}\n\n\n\n")
                marker._parent._parent.left += int(value * _MANAGEMENT_SUMMARY_MARKER_RANGE)


class ModernizationVolumeMarkerPlaceholder(_ManagementSummaryMarkerPlaceholder):
    key = "MARKER_MODERNIZATION_VOLUME"

    @classmethod
    def value(cls, parameter=None) -> tuple[float, str]:
        return modernization_data.total_volume / 5000.0, f"{round(modernization_data.total_volume)} PY"


class ModernizationTechnicalDebtMarkerPlaceholder(_ManagementSummaryMarkerPlaceholder):
    key = "MARKER_MODERNIZATION_TECHNICAL_DEBT"

    @classmethod
    def value(cls, parameter=None) -> tuple[float, str]:
        technical_debt = sum(
            candidate.technical_debt_in_py for candidate in modernization_data.modernization_candidates)
        return technical_debt / modernization_data.total_volume, f"{round(technical_debt)} PY"


class ModernizationSpeedMarkerPlaceholder(_ManagementSummaryMarkerPlaceholder):
    key = "MARKER_MODERNIZATION_SPEED"

    @classmethod
    def value(cls, parameter=None) -> tuple[float, str]:
        total = 0.0
        total_weight = 0.0

        for candidate in modernization_data.modernization_candidates:
            total += candidate.estimated_change_speed * candidate.volume_in_py
            total_weight += candidate.volume_in_py

        weighted_average = (total / total_weight) / 100.0 if total_weight > 0.0 else 0.0
        return weighted_average, f"+ {round(weighted_average * 100)}%"


class ModernizationEffortMarkerPlaceholder(_ManagementSummaryMarkerPlaceholder):
    key = "MARKER_MODERNIZATION_EFFORT"

    @classmethod
    def value(cls, parameter=None) -> tuple[float, str]:
        effort = sum(candidate.estimated_effort_py for candidate in modernization_data.modernization_candidates)
        return effort / modernization_data.total_volume, f"{round(effort)} PY"
