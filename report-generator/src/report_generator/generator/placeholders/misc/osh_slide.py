#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from typing import Callable, Tuple

from pptx.chart.data import ChartData
from pptx.presentation import Presentation
from pptx.slide import Slide

from report_generator.generator import report_utils
from report_generator.generator.data_models import osh_data
from report_generator.generator.placeholders import Placeholder
from report_generator.generator.placeholders.base import PlaceholderDocType


def _format_chart_data(data) -> Tuple[ChartData, ChartData]:
    chart_data = ChartData()
    chart_data.categories = ["Vulnerability risk", "Legal risk"]
    chart_data.add_series("Critical risk", [data.vuln_risks[0], data.license_risks[0]])
    chart_data.add_series("High risk", [data.vuln_risks[1], data.license_risks[1]])
    chart_data.add_series("Medium risk", [data.vuln_risks[2], data.license_risks[2]])
    chart_data.add_series("Low risk", [data.vuln_risks[3], data.license_risks[3]])

    chart_data2 = ChartData()
    chart_data2.categories = ["Freshness risk", "Stability risk", "Management risk", "Activity risk"]
    chart_data2.add_series("Critical risk", [data.freshness_risks[0], data.stability_risks[0], data.mgmt_risks[0],
                                             data.activity_risks[0]])
    chart_data2.add_series("High risk", [data.freshness_risks[1], data.stability_risks[1], data.mgmt_risks[1],
                                         data.activity_risks[1]])
    chart_data2.add_series("Medium risk", [data.freshness_risks[2], data.stability_risks[2], data.mgmt_risks[2],
                                           data.activity_risks[2]])
    chart_data2.add_series("Low risk", [data.freshness_risks[3], data.stability_risks[3], data.mgmt_risks[3],
                                        data.activity_risks[3]])

    return chart_data, chart_data2


def _determine_chart_axis_max():
    data = osh_data.data  # TODO actually use chart data
    max_bar_length = max(sum(data.vuln_risks[0:4]), sum(data.license_risks[0:4]), sum(data.freshness_risks[0:4]),
                         sum(data.activity_risks[0:4]), sum(data.stability_risks[0:4]), sum(data.mgmt_risks[0:4]))

    return max_bar_length * 1.1


def _populate_osh_system_slide(slide: Slide, data, data2):
    shapes_by_name = dict((s.name, s) for s in slide.shapes)
    chart_axis_max = _determine_chart_axis_max()
    _set_chart_data_and_axis(shapes_by_name["CHART_1"].chart, data, chart_axis_max)
    _set_chart_data_and_axis(shapes_by_name["CHART_2"].chart, data2, chart_axis_max)


def _set_chart_data_and_axis(chart, data, axis_max):
    chart.replace_data(data)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = axis_max


class OSHSlidePlaceholder(Placeholder):
    """Traditional SIG OSH system-level slide, with risk bar charts for all 6 OSH metrics."""
    key = "OSH_SLIDE"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _format_chart_data(osh_data.data)

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        slides = report_utils.pptx.identify_specific_slide(presentation, key)

        if len(slides) == 0:
            return

        (data, data2) = value_cb()
        for slide in slides:
            _populate_osh_system_slide(slide, data, data2)
