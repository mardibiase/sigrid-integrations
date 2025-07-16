#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import logging
import re
from typing import Union

from pptx import Presentation
from pptx.dml.color import ColorFormat, RGBColor, _Color
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.table import Table, _Cell
from pptx.text.text import _Paragraph, _Run

NA_STAR_COLOR = RGBColor(0xb5, 0xb5, 0xb5)
ONE_STAR_COLOR = RGBColor(0xdb, 0x4a, 0x3d)
TWO_STAR_COLOR = RGBColor(0xef, 0x98, 0x1a)
THREE_STAR_COLOR = RGBColor(0xf8, 0xc6, 0x40)
FOUR_STAR_COLOR = RGBColor(0x57, 0xc9, 0x68)
FIVE_STAR_COLOR = RGBColor(0x2c, 0x96, 0x3f)


def print_slide_ids(slide):
    # Print slide IDs and names for debugging purposes
    logging.debug("Placeholders:")
    for shape in slide.placeholders:
        logging.debug('%d %s' % (shape.placeholder_format.idx, shape.name))
    logging.debug("----\n")
    logging.debug("Shapes:")
    for shape in slide.shapes:
        logging.debug('%d [%s] %s' % (shape.shape_id, shape.name, "(This is a chart)" if shape.has_chart else ""))


def update_many_paragraphs(paragraphs, placeholder_id, replacement_text, font={}):
    for paragraph in paragraphs:
        update_paragraph(paragraph, placeholder_id, replacement_text, font)


def update_paragraph(paragraph, placeholder_id, replacement_text, font={}):
    if paragraph:
        # Powerpoint sometimes puts arbitrarily splits up text in runs, even if they have the same formatting.
        # This sometimes puts one of our placeholders into 2 or more runs (e.g. AAP_NOOT_MIES might be be "AAP_", "NOOT", "_MIES")
        # Here we stitch those back together so that we can do effective repacement
        merge_runs_with_same_formatting(paragraph)

        run_with_placeholder = None
        if paragraph.runs:
            for run in paragraph.runs:
                if placeholder_id in run.text:
                    run_with_placeholder = run
                    break

        if run_with_placeholder:
            run_with_placeholder.text = run_with_placeholder.text.replace(placeholder_id, str(replacement_text))
            logging.debug(
                f"Replacing: {placeholder_id} with \"{replacement_text}\". New text: {run_with_placeholder.text}")
            update_run_font(run_with_placeholder, font)


def find_shapes_with_text(presentation, search_text):
    shapes = []
    for slide in presentation.slides:
        paragraphs = find_text_in_slide(slide, search_text)
        # A paragraph is typically in a TextGroup which is in a Shape, so we call getparent() twice
        shapes += [paragraph._parent._parent for paragraph in paragraphs]
    return shapes


def find_text_in_presentation(presentation, search_text):
    paragraphs = []
    for slide in presentation.slides:
        paragraphs.extend(find_text_in_slide(slide, search_text))

    return paragraphs


def find_text_in_slide(slide, search_text):
    paragraphs = []
    for shape in slide.shapes:
        result = find_text_in_shape(shape, search_text)
        if result:
            paragraphs.append(result)
    return paragraphs


def find_text_in_shape(shape, search_text):
    # If the shape is a table, iterate through its cells
    if "GraphicFrame" in type(shape).__name__:
        if shape.has_table:
            for cell in shape.table.iter_cells():
                if re.match(fr".*\b{search_text}\b.*", cell.text):
                    return cell.text_frame.paragraphs[0]
        else:
            return None

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if re.match(fr".*\b{search_text}\b.*", paragraph.text):
                return paragraph

    # If the shape is a group, iterate through its shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            result = find_text_in_shape(s, search_text)
            if result:
                return result

    return None


def add_content_paragraph(self, text_frame, markers, content, paragraph=None):
    if paragraph is None:
        paragraph = text_frame.add_paragraph()
    for marker in markers:
        set_sig_marker(paragraph, marker)
    run = paragraph.add_run()
    run.text = " " + content


def set_sig_marker(paragraph, marker):
    run = paragraph.add_run()
    run.text = marker
    run.font.name = "SIGMarker"

    # Red, yellow and green colors are taken from the SIG pptx template Signal colors
    if marker == "-":
        run.font.color.rgb = RGBColor(0xcb, 0x55, 0x45)
    if marker == "=":
        run.font.color.rgb = RGBColor(0xf0, 0xc8, 0x5a)
    if marker == "+":
        run.font.color.rgb = RGBColor(0x77, 0xc6, 0x73)


def update_run_font(run, font):
    if font:
        for k, v in font.items():
            if k == "bold" and v:
                run.font.bold = True
            if k == "name":
                run.font.name = v
            if k == "size":
                run.font.size = v


def merge_runs_with_same_formatting(paragraph):
    last_run = None
    for run in paragraph.runs:
        if last_run is None:
            last_run = run
            continue
        if has_same_formatting(run, last_run):
            last_run = combine_runs(last_run, run)
            continue
        last_run = run


def has_same_formatting(run, run_2):
    font, font_2 = run.font, run_2.font
    if font.bold != font_2.bold:
        return False
    if font.italic != font_2.italic:
        return False
    if font.name != font_2.name:
        return False
    if font.size != font_2.size:
        return False
    if font.underline != font_2.underline:
        return False
    return True


def combine_runs(base, suffix):
    base.text = base.text + suffix.text
    r_to_remove = suffix._r
    r_to_remove.getparent().remove(r_to_remove)
    return base


def add_xml_element(parent_xml, tag, **attrs):
    element = OxmlElement(tag)
    element.attrib.update(attrs)
    parent_xml.append(element)
    return element


# inspired by https://groups.google.com/g/python-pptx/c/UTkdemIZICw
# TODO: debug the powerpoint and make this work
def set_cell_border(cell, border_color, border_width, border_style, *border_types):
    logging.debug(f'setting borders for {border_types}, {border_color=}, {border_style=}, {border_width=}')
    cell_xml = cell._tc
    cell_xml_parent = cell_xml.get_or_add_tcPr()

    def fill_line(line, border_color, border_style):
        line_fill = add_xml_element(line, 'a:solidFill')
        line_rgb = add_xml_element(line_fill, 'a:srgbClr', val=border_color)
        line_style = add_xml_element(line, 'a:prstDash', val=border_style)
        line_round_ = add_xml_element(line, 'a:round')
        line_headEnd = add_xml_element(line, 'a:headEnd', type='none', w='med', len='med')
        line_tailEnd = add_xml_element(line, 'a:tailEnd', type='none', w='med', len='med')

    border_type_tag = {
        'left'  : 'a:lnL',
        'right' : 'a:lnR',
        'top'   : 'a:lnT',
        'bottom': 'a:lnB'
    }

    for border_type in border_types:
        tag = border_type_tag[border_type]
        line = add_xml_element(cell_xml_parent, tag, w=str(border_width), cap='flat', cmpd='sng',
                               algn='ctr')
        fill_line(line, border_color, border_style)


def set_shape_color(shape, rgbColor):
    shape.fill.fore_color.rgb = rgbColor


def identify_specific_slide(presentation, marker):
    specific_slides = []
    for slide in presentation.slides:
        if find_text_in_slide(slide, marker):
            specific_slides.append(slide)
    return specific_slides


def determine_rating_color(rating):
    if rating < 0.1:
        return NA_STAR_COLOR
    if rating < 1.5:
        return ONE_STAR_COLOR
    elif rating < 2.5:
        return TWO_STAR_COLOR
    elif rating < 3.5:
        return THREE_STAR_COLOR
    elif rating < 4.5:
        return FOUR_STAR_COLOR
    else:
        return FIVE_STAR_COLOR


def test_code_ratio_color(ratio):
    if ratio <= 0.01:
        return ONE_STAR_COLOR
    elif ratio <= 0.15:
        return TWO_STAR_COLOR
    elif ratio <= 0.5:
        return THREE_STAR_COLOR
    elif ratio <= 1.5:
        return FOUR_STAR_COLOR
    else:
        return FIVE_STAR_COLOR


def gather_charts(presentation: Presentation, key: str):
    charts = []
    for slide in identify_specific_slide(presentation, key):
        for shape in slide.shapes:
            if shape.has_chart:
                charts.append(shape.chart)
    return charts


def find_tables(presentation: Presentation, key: str):
    return [
        shape.table
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table and shape.name == key
    ]

"""
Fills a table with the provided values.
- The table must have the same number of rows and columns as the value list. @@ make this more flexible
- The table will override existing content in the cells.
- If a cell has existing formatting, that formatting will be applied to all consecutive rows in that column.
- If there is no cell for a column that defines the formatting, some default formatting will be used, which is likely not what you want.
"""
def update_table(table: Table, value: list[list[Union[str, int, float]]]):
    # Dictionary to store reference runs for each column
    column_fonts = {}

    # Second pass: update the table
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            if row_idx >= len(value) or col_idx >= len(value[row_idx]):
                ## @@TODO add row or column if not present
                return

            paragraph: _Paragraph = cell.text_frame.paragraphs[0]
            if paragraph.runs:
                column_fonts[col_idx] = _copy_font_properties(paragraph.runs[0])
            paragraph.clear()

            new_text = str(value[row_idx][col_idx])

            run: _Run = paragraph.add_run()
            run.text = new_text

            if column_fonts[col_idx]:
                _apply_font_properties(run, column_fonts[col_idx])


def _copy_font_properties(source_run: _Run):
    font = source_run.font
    color: ColorFormat = font.color
    return {
        'bold': font.bold,
        'italic': font.italic,
        'name': font.name,
        'size': font.size,
        'underline': font.underline,
        'color': {
            'rgb': color.rgb if hasattr(color, 'rgb') else None,
            'theme_color': color.theme_color if hasattr(color, 'theme_color') else None,
            'brightness': color.brightness if hasattr(color, 'brightness') else None,
        }
    }

def _apply_font_properties(target_run: _Run, font_properties: dict):
    target_run.font.bold = font_properties['bold']
    target_run.font.italic = font_properties['italic']
    target_run.font.name = font_properties['name']
    target_run.font.size = font_properties['size']
    target_run.font.underline = font_properties['underline']

    if font_properties['color']['rgb'] is not None:
        target_run.font.color.rgb = font_properties['color']['rgb']
    if font_properties['color']['theme_color'] is not None and font_properties['color']['theme_color'] is not MSO_THEME_COLOR.NOT_THEME_COLOR:
        print(font_properties['color']['theme_color'])
        target_run.font.color.theme_color = font_properties['color']['theme_color']
    if font_properties['color']['brightness'] is not None and target_run.font.color.type is not None:
        target_run.font.color.brightness = font_properties['color']['brightness']