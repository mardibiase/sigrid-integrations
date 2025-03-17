from pptx import Presentation
from pptx.chart.data import XyChartData

from report_generator import report_utils
from report_generator.data_models import maintainability_data
from report_generator.placeholders import Placeholder


class MaintainabilityGalaxyChartPlaceholder(Placeholder):
    key = "GALAXY_SLIDE"

    @classmethod
    def value(cls, parameter=None):
        pass

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb) -> None:
        charts = []

        for slide in report_utils.pptx.identify_specific_slide(presentation, key):
            shapes_by_name = dict((s.name, s) for s in slide.shapes)
            # Todo autodetect galaxy chart index. (Or just assume it's the first and only chart on this slide, see tkovac generator implementation)
            charts.append(shapes_by_name["CHART_1"].chart)

        if len(charts) == 0:
            return

        volume = maintainability_data.system_py
        maint_rating = maintainability_data.maintainability_rating
        system_name = maintainability_data.system_name

        chart_data = XyChartData()
        series = chart_data.add_series("Series 1")

        series.add_data_point(volume, maint_rating)

        for chart in charts:
            chart.replace_data(chart_data)
            chart.series[0].points[0].data_label.text_frame.text = system_name
