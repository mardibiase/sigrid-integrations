# Copyright Software Improvement Group
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
# 
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import pptx
import logging
from .report import Report

from .osh_report import OSHReport, OSHData
from .slide_utils import SlideUtils

from pptx.chart.data import XyChartData, CategoryChartData
from pptx.dml.color import RGBColor

# Use absolute imports so that test code can also work
from .report import Report

class PowerpointReport(Report):

    # This is the number of "EMU" we need to move the slide marker for every 1.0 change in rating
    # in our current slide template
    # Value determined through experimentation. Assuming the marker starts at 3.0
    RATING_MARKER_MOVE_SIZE = 2200000

    NA_STAR_COLOR = RGBColor(0xb5, 0xb5, 0xb5)
    ONE_STAR_COLOR = RGBColor(0xdb, 0x4a, 0x3d)
    TWO_STAR_COLOR = RGBColor(0xef, 0x98, 0x1a)
    THREE_STAR_COLOR = RGBColor(0xf8, 0xc6, 0x40)
    FOUR_STAR_COLOR = RGBColor(0x57, 0xc9, 0x68)
    FIVE_STAR_COLOR = RGBColor(0x2c, 0x96, 0x3f)

    def __init__(self, customer, system, token, out_file, template):
        super().__init__(customer, system, token)
        self.out_file = out_file
        self.presentation = pptx.Presentation(template)
    
    def save_file(self):
        logging.info(f"Written output to {self.out_file}")
        self.presentation.save(self.out_file)

    def create_osh_report_specific(self, osh_data: OSHData, osh_report: OSHReport):
        slides = self.identify_specific_slide("OSH_SLIDE")
        osh_report.populate_charts(slides, osh_data, False)
        
    def fill_metric_report_specific(self, model, metric, rating):
        self.color_shape_by_rating("COLOR_" + model + "_RATING_" + metric, rating)

    def color_shape_by_rating(self, placeholder_text, rating):
        rating_color = self.determine_rating_color(rating)
        for shape in SlideUtils.find_shapes_with_text(self.presentation, placeholder_text):
            SlideUtils.set_shape_color(shape, rating_color)
    
    def determine_rating_color(self, rating):
        if rating < 0.1:
            return self.NA_STAR_COLOR
        if rating < 1.5:
            return self.ONE_STAR_COLOR
        elif rating < 2.5:
            return self.TWO_STAR_COLOR
        elif rating < 3.5:
            return self.THREE_STAR_COLOR
        elif rating < 4.5:
            return self.FOUR_STAR_COLOR
        else:
            return self.FIVE_STAR_COLOR

    def create_maintainability_report_specific(self, maint_data):

        self.generate_maintainability_galaxy_chart(maint_data["system"].capitalize(), maint_data["maintainability"], maint_data["volumeInPersonMonths"] / 12)

        self.update_movable_marker("MARKER_MAINT_RATING", maint_data["maintainability"])

        #TODO: Add after maintainability API changes are merged, remove this mock data
        #maint_data["technologies"] = [PowerpointReport._mock_technology("aap", 1, 15, 3.0, 0.02, "TARGET"),
        #        PowerpointReport._mock_technology("noot", 3, 2, 3.0, 0.15, "TARGET"),
        #        PowerpointReport._mock_technology("mies", 2, 87, 4.0, 1.3, "TARGET"),
        #        PowerpointReport._mock_technology("wim", 15, 87, 3.0, 1.7, "TARGET"),
        #        PowerpointReport._mock_technology("zus", 7, 87, 3.0, 15.0, "TARGET"),
        #        PowerpointReport._mock_technology("jet", 1, 87, 3.0, 0.4, "TOLERATE"),
        #        PowerpointReport._mock_technology("teun", 18, 87, 3.0, 0.9, "TARGET"),]
        #self.fill_technologies(maint_data["technologies"])
    
    def create_technology_report_specific(self, unsorted_tech_data, sorted_tech_data):
    #def fill_technologies(self, tech_data):

        total_volume_pm = sum([data["volumeInPersonMonths"] for data in sorted_tech_data])

        self.generate_category_chart("TECHNOLOGY_CHART",
                                 [data["displayName"] for data in sorted_tech_data],
                                 [[data["volumeInPersonMonths"] / total_volume_pm for data in sorted_tech_data]])
        
        self.generate_category_chart("TEST_CODE_RATIO_CHART", 
                                 [data["displayName"] for data in sorted_tech_data],
                                 [[data["volumeInPersonMonths"] / total_volume_pm for data in sorted_tech_data]],
                                 [PowerpointReport.test_code_ratio_color(data["testCodeRatio"]) for data in sorted_tech_data])
    
    def target_ratio_for_technology(risk_level, volume_percentage):
        if risk_level == "TARGET":
            return volume_percentage
        elif risk_level == "TOLERATE":
            return 0.5 * volume_percentage
        else:
            return 0
        
    def test_code_ratio_color(ratio):
        if ratio <= 0.01:
            return PowerpointReport.ONE_STAR_COLOR
        elif ratio <= 0.15:
            return PowerpointReport.TWO_STAR_COLOR
        elif ratio <= 0.5:
            return PowerpointReport.THREE_STAR_COLOR
        elif ratio <= 1.5:
            return PowerpointReport.FOUR_STAR_COLOR
        else:
            return PowerpointReport.FIVE_STAR_COLOR
        
    def generate_category_chart(self, slide_marker, x_values, y_series, colors = None):
        for slide in self.identify_specific_slide(slide_marker):
            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    chart_data = CategoryChartData()
                    chart_data.categories = x_values
                    for y in y_series:
                        chart_data.add_series("Volume in Person Months", y)
                    chart.replace_data(chart_data)

                    if colors != None:
                        for serie in chart.series:
                            for idx, point in enumerate(serie.points):
                                point.format.fill.solid()
                                point.format.fill.fore_color.rgb = colors[idx]



    def update_movable_marker(self, placeholder, rating):
        rating = Report.maintainability_round(rating)
        
        for slide in self.presentation.slides:
            paragraphs = SlideUtils.find_text_in_slide(slide, placeholder)
            for paragraph in paragraphs:
                marker = paragraph._parent._parent._parent._parent
                marker.left = int(marker.left + self.RATING_MARKER_MOVE_SIZE * self.distance_to_average(rating))

        self.update_placeholder(placeholder, rating)
    
    # e.g. 0.5 returns -2.5. 4.0 returns +1.0
    def distance_to_average(self, rating):
        return -1*(3.0 - rating)

    def generate_maintainability_galaxy_chart(self, system_name, maint_rating, volumeInPy):
        for slide in self.identify_specific_slide("GALAXY_SLIDE"):
            shapes_by_name = dict((s.name, s) for s in slide.shapes)
            # Todo autodetect galaxy chart index. (Or just assume it's the first and only chart on this slide, see tkovac generator implementation)
            chart1 = shapes_by_name["CHART_1"].chart
            chart_data = XyChartData()
            series = chart_data.add_series("Series 1")
            
            series.add_data_point(volumeInPy, maint_rating)
            chart1.replace_data(chart_data)
            chart1.series[0].points[0].data_label.text_frame.text = system_name

    def create_architecture_quality_report_specific(self, architecture_data):
       ## for metric in StaticData.arch_metrics:
        #    self.fill_metric("ARCH", metric, architecture_data["ratings"]["systemProperties"][GenericUtils.to_json_name(metric)])
        
       # for metric in self.arch_subcharacteristics:
        #    self.fill_metric("ARCH", metric, architecture_data["ratings"]["subcharacteristics"][GenericUtils.to_json_name(metric)])
            
        self.update_movable_marker("MARKER_ARCH_RATING", architecture_data["ratings"]["architecture"])
        
    
    def update_placeholder(self, placeholder_id, replacement_text, font={}):
        for slide in self.presentation.slides:
            self.update_slide(slide, placeholder_id, replacement_text, font)

    def update_slide(self, slide, placeholder_id, replacement_text, font={}):
        paragraphs = SlideUtils.find_text_in_slide(slide, placeholder_id)
        for paragraph in paragraphs:
            SlideUtils.update_paragraph(paragraph, placeholder_id, replacement_text, font)

    def identify_specific_slide(self, marker):
        specific_slides = []
        for slide in self.presentation.slides:
            if SlideUtils.find_text_in_slide(slide, marker):
                specific_slides.append(slide)
        return specific_slides


        
