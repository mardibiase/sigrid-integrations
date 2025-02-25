# Copyright Software Improvement Group
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
# 
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import dateutil.parser
from pptx.chart.data import ChartData
from .slide_utils import SlideUtils

class OSHVuln:
    name = ""
    version = ""
    cve = ""
    severity = 0.0

    def __init__(self, name, version, cve, severity):
        self.name = name
        self.version = version
        self.cve = cve
        self.severity = severity

class OSHData:
    total_deps = 0

    date_day = ""
    date_month = ""
    date_year = ""

    # critical, high, medium, low, no risk
    vuln_risks = [0,0,0,0,0]
    license_risks = [0,0,0,0,0]
    freshness_risks = [0,0,0,0,0]
    stability_risks = [0,0,0,0,0]
    mgmt_risks = [0,0,0,0,0]
    activity_risks = [0,0,0,0,0]

    vulns = []

    def total_vulnerable(self):
        return sum(self.vuln_risks[0:4]) # to add all columns for this metric

class OSHReport:
    
    def aggregate_data(self, input_data):
        data = OSHData()        
        for component in input_data.get("components", []):
            data.total_deps += 1
            if(data.date_year == ""):
                (data.date_year, data.date_month, data.date_day) = OSHReport.format_date(input_data["metadata"]["timestamp"])

            OSHReport.assign_risk(data.vuln_risks, OSHReport.find_cyclonedx_property_value(component["properties"], "sigrid:risk:vulnerability"))
            OSHReport.assign_risk(data.license_risks, OSHReport.find_cyclonedx_property_value(component["properties"], "sigrid:risk:legal"))
            OSHReport.assign_risk(data.freshness_risks, OSHReport.find_cyclonedx_property_value(component["properties"], "sigrid:risk:freshness"))
            OSHReport.assign_risk(data.stability_risks, OSHReport.find_cyclonedx_property_value(component["properties"], "sigrid:risk:stability"))
            OSHReport.assign_risk(data.mgmt_risks, OSHReport.find_cyclonedx_property_value(component["properties"], "sigrid:risk:management"))
            OSHReport.assign_risk(data.activity_risks, OSHReport.find_cyclonedx_property_value(component["properties"], "sigrid:risk:activity"))
        return data
    
    def populate_charts(self, slides, osh_data: OSHData, print_ids):
        for slide in slides:
            OSHReport.populate_osh_system_slide(slide, osh_data, print_ids)
        
    def populate_osh_system_slide(slide, osh_data: OSHData, print_ids=False):

        if(print_ids):
            SlideUtils.print_slide_ids(slide)

        # For easy access to all shapes and placeholders, we store them in a dictionary by their name.
        shapes_by_name = dict((s.name, s) for s in slide.shapes)
        placeholders_by_name = dict((p.name, p) for p in slide.placeholders)

        (data, data2) = OSHReport.format_chart_data(osh_data)
        chart_axis_max = OSHReport.determine_chart_axis_max(osh_data)
        # Todo: Consider finding these chart names dynamically instead of hardcoding them. CHallenge is we have 2 charts on one slide.
        OSHReport.set_chart_data_and_axis(shapes_by_name["CHART_1"].chart, data, chart_axis_max)
        OSHReport.set_chart_data_and_axis(shapes_by_name["CHART_2"].chart, data2, chart_axis_max)

    def format_chart_data(osh_data: OSHData): 
        data = ChartData()
        data.categories = ["Vulnerability risk", "Legal risk"]
        data.add_series("Critical risk", [osh_data.vuln_risks[0], osh_data.license_risks[0]])
        data.add_series("High risk", [osh_data.vuln_risks[1], osh_data.license_risks[1]])
        data.add_series("Medium risk", [osh_data.vuln_risks[2], osh_data.license_risks[2]])
        data.add_series("Low risk", [osh_data.vuln_risks[3], osh_data.license_risks[3]])
        data2 = ChartData()
        data2.categories = ["Freshness risk", "Stability risk", "Management risk", "Activity risk"]
        data2.add_series("Critical risk", [osh_data.freshness_risks[0], osh_data.stability_risks[0], osh_data.mgmt_risks[0], osh_data.activity_risks[0]])
        data2.add_series("High risk", [osh_data.freshness_risks[1], osh_data.stability_risks[1], osh_data.mgmt_risks[1], osh_data.activity_risks[1]])
        data2.add_series("Medium risk", [osh_data.freshness_risks[2], osh_data.stability_risks[2], osh_data.mgmt_risks[2], osh_data.activity_risks[2]])
        data2.add_series("Low risk", [osh_data.freshness_risks[3], osh_data.stability_risks[3], osh_data.mgmt_risks[3], osh_data.activity_risks[3]])

        return (data, data2)
    
    def set_chart_data_and_axis(chart, data, max):
        chart.replace_data(data)
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = max

    def determine_chart_axis_max(osh_data: OSHData):
        max_bar_length = max(sum(osh_data.vuln_risks[0:4]), sum(osh_data.license_risks[0:4]), sum(osh_data.freshness_risks[0:4]), sum(osh_data.activity_risks[0:4]), sum(osh_data.stability_risks[0:4]), sum(osh_data.mgmt_risks[0:4]))
        
        # add ±10% of padding
        return max_bar_length * 1.1
    
    def vulnerability_summary(osh_data: OSHData):
        total_vulns = sum(osh_data.vuln_risks[0:4])
        if total_vulns > 0:
            pct_vulns = max(total_vulns / osh_data.total_deps, 0.01) # Percentage should always be at least 1. 0% looks stupid.
            return f"{pct_vulns:.0%} of dependencies ({total_vulns} in total) used in the system contain one or more known vulnerabilities."
        else:
            return "The system is free of known vulnerabilities."
        
    def freshness_summary(osh_data: OSHData):
        total_outdated = sum(osh_data.freshness_risks[0:3]) # Only count critial+high+medium risk. Llow is fresh enough to not report on
        if total_outdated > 0:
            pct_outdated = max(total_outdated / osh_data.total_deps, 0.01)
            return f"{pct_outdated:.0%} of dependencies ({total_outdated} in total) used in the system have not been updated for over 2 years."
        else:
            return "All dependencies in the system have been updated in the last 2 years."
    
    def legal_summary(osh_data: OSHData):
        total_legal = sum(osh_data.license_risks[0:3]) # Only count critial, high, and medium. Low license risk is typically not restrictive, so not interesting to report on
        if total_legal > 0:
            pct_legal = max(total_legal / osh_data.total_deps, 0.01)
            return f"{pct_legal:.0%} of dependencies ({total_legal} in total) uses a potentially restrictive open-source license (e.g. GPL/AGPL)."
        else:
            return "All dependencies in the system use relatively liberal open-source licenses."
        
    def management_summary(osh_data: OSHData):
        total_unmanaged = sum(osh_data.mgmt_risks[0:4])
        if total_unmanaged > 0:
            pct_unmanaged = max(total_unmanaged / osh_data.total_deps, 0.01)
            return f"{pct_unmanaged:.0%} of dependencies ({total_unmanaged} in total) does not use a package manager but is placed in the codebase directly."
        else:
            return "All dependencies in the sysstem are managed by a package manager."
        
    def format_date(input_date):
        date = dateutil.parser.isoparse(input_date)
        return (str(date.year), date.strftime('%b').upper(), str(date.day))

    def find_cyclonedx_property_value(properties, key):
        for property in properties:
            if property["name"] == key:
                return property["value"]
        return None
    
    def assign_risk(values, risk):
        if(risk == "CRITICAL"):
            values[0] += 1
        elif(risk == "HIGH"):
            values[1] += 1
        elif(risk == "MEDIUM"):
            values[2] += 1
        elif(risk == "LOW"):
            values[3] += 1
        else:
            values[4] += 1
